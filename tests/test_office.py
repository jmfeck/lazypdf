import os
import shutil
import sys

import pytest

import lazypdf as lz

has_msoffice = False
if sys.platform == "win32":
    try:
        import win32com.client

        # Try to instantiate Word to confirm Office is actually installed
        win32com.client.Dispatch("Word.Application").Quit()
        has_msoffice = True
    except Exception:
        pass
has_libreoffice = shutil.which("soffice") is not None or shutil.which("libreoffice") is not None
has_office = has_msoffice or has_libreoffice

skip_no_office = pytest.mark.skipif(not has_office, reason="No Office suite available")


@pytest.fixture
def sample_docx(tmp_path):
    """Create a simple .docx file for testing."""
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not installed")

    path = str(tmp_path / "sample.docx")
    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is a test paragraph.")
    doc.add_paragraph("Second paragraph with more text.")
    doc.save(path)
    return path


@pytest.fixture
def sample_xlsx(tmp_path):
    """Create a simple .xlsx file for testing."""
    try:
        from openpyxl import Workbook
    except ImportError:
        pytest.skip("openpyxl not installed")

    path = str(tmp_path / "sample.xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Age", "City"])
    ws.append(["Alice", 30, "Sao Paulo"])
    ws.append(["Bob", 25, "Porto Alegre"])
    wb.save(path)
    return path


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a simple .pptx file for testing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches  # noqa: F401
    except ImportError:
        pytest.skip("python-pptx not installed")

    path = str(tmp_path / "sample.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "This is a test presentation."
    prs.save(path)
    return path


@pytest.fixture
def sample_csv(tmp_path):
    """Create a simple .csv file for testing."""
    path = str(tmp_path / "sample.csv")
    with open(path, "w") as f:
        f.write("Name,Age,City\n")
        f.write("Alice,30,Sao Paulo\n")
        f.write("Bob,25,Porto Alegre\n")
    return path


@skip_no_office
class TestReadDocx:
    def test_read_docx(self, sample_docx, output_dir):
        pdf = lz.read_docx(sample_docx)
        assert len(pdf) >= 1
        text = pdf.extract_text()
        assert "Test Document" in text
        pdf.close()

    def test_read_docx_to_pdf(self, sample_docx, output_dir):
        out = lz.read_docx(sample_docx).to_pdf(f"{output_dir}/from_docx.pdf")
        assert os.path.exists(out)

    def test_read_docx_chain(self, sample_docx, output_dir):
        out = (
            lz.read_docx(sample_docx)
            .add_watermark("DRAFT")
            .add_page_numbers()
            .compress()
            .to_pdf(f"{output_dir}/docx_chained.pdf")
        )
        assert os.path.exists(out)


@skip_no_office
class TestReadXlsx:
    def test_read_xlsx(self, sample_xlsx):
        pdf = lz.read_xlsx(sample_xlsx)
        assert len(pdf) >= 1
        pdf.close()

    def test_read_xlsx_to_pdf(self, sample_xlsx, output_dir):
        out = lz.read_xlsx(sample_xlsx).to_pdf(f"{output_dir}/from_xlsx.pdf")
        assert os.path.exists(out)

    def test_read_xlsx_to_images(self, sample_xlsx, output_dir):
        paths = lz.read_xlsx(sample_xlsx).to_png(output_dir)
        assert len(paths) >= 1


@skip_no_office
class TestReadPptx:
    def test_read_pptx(self, sample_pptx):
        pdf = lz.read_pptx(sample_pptx)
        assert len(pdf) >= 1
        pdf.close()

    def test_read_pptx_to_pdf(self, sample_pptx, output_dir):
        out = lz.read_pptx(sample_pptx).to_pdf(f"{output_dir}/from_pptx.pdf")
        assert os.path.exists(out)

    def test_read_pptx_chain(self, sample_pptx, output_dir):
        out = lz.read_pptx(sample_pptx).rotate(90).to_pdf(f"{output_dir}/pptx_rotated.pdf")
        assert os.path.exists(out)


@skip_no_office
class TestReadCsv:
    def test_read_csv(self, sample_csv):
        pdf = lz.read_csv(sample_csv)
        assert len(pdf) >= 1
        pdf.close()

    def test_read_csv_to_pdf(self, sample_csv, output_dir):
        out = lz.read_csv(sample_csv).to_pdf(f"{output_dir}/from_csv.pdf")
        assert os.path.exists(out)
